from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parent.parent / "docs" / "Wisdom-Weasel_演示版PPT_新版_v3.pptx"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
ORANGE = RGBColor(245, 158, 11)
AMBER = RGBColor(251, 191, 36)
CYAN = RGBColor(34, 211, 238)
RED = RGBColor(239, 68, 68)
GREEN = RGBColor(16, 185, 129)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 250, 252)
SOFT = RGBColor(241, 245, 249)
LINE = RGBColor(226, 232, 240)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def top_band(slide, text, accent=ORANGE):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.33)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.03), Inches(5.5), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE


def title_block(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.82), Inches(0.72), Inches(11.7), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.85), Inches(1.38), Inches(11.0), Inches(0.4))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED


def footer(slide, text="Wisdom-Weasel / 演示版"):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(7.05), Inches(11.8), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.82), Inches(7.08), Inches(5.5), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def card(slide, left, top, width, height, title, lines, accent=ORANGE, title_size=18, body_size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.14), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(left + Inches(0.26), top + Inches(0.16), width - Inches(0.4), height - Inches(0.2))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.font.name = "Microsoft YaHei"
        para.font.size = Pt(body_size)
        para.font.color.rgb = SLATE
        para.space_before = Pt(3)


def bullets(slide, left, top, width, height, items, size=18, color=SLATE):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def callout(slide, left, top, width, height, text, fill=SOFT, line=LINE, color=NAVY, size=16):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tx = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), width - Inches(0.35), height - Inches(0.18))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color


def arch_box(slide, left, top, width, height, title, body, fill, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tx = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.12), width - Inches(0.3), height - Inches(0.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = text_color
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.font.color.rgb = text_color


def arrow(slide, left, top):
    tx = slide.shapes.add_textbox(left, top, Inches(0.55), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = "→"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORANGE


def add_table_slide_rows(slide, left, top, width, height, headers, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    col_widths = [1.35, 3.2, 0.8, 3.5, 3.5]
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = Inches(w)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else SOFT
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(1.05), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = ORANGE
    side.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.78), 0, Inches(0.16), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.42), Inches(1.1), Inches(10.8), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = "Wisdom-Weasel"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(1.45), Inches(2.02), Inches(10.6), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "基于上下文感知候选重排与联想预测的智能中文输入法"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = SLATE
    callout(
        slide,
        Inches(1.45),
        Inches(3.0),
        Inches(10.0),
        Inches(0.78),
        "让输入法从“拼音转文字工具”升级为“表达协作系统”",
        fill=SOFT,
        line=ORANGE,
        color=ORANGE,
        size=18,
    )
    bullets(
        slide,
        Inches(1.5),
        Inches(4.2),
        Inches(10.3),
        Inches(1.5),
        [
            "解决的不是“能不能打出字”，而是“能不能更顺着人的思路完成表达”",
            "核心方向：理解语境、优化候选顺序、在停顿时继续辅助表达",
        ],
        size=18,
    )
    footer(slide, "Wisdom-Weasel / 演示版")

    # 2 problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band(slide, "01 问题定义", ORANGE)
    title_block(slide, "传统输入法真正的问题：不理解“我现在在写什么”")
    bullets(
        slide,
        Inches(0.95),
        Inches(1.95),
        Inches(11.4),
        Inches(1.8),
        [
            "同一组拼音在不同语境下，用户真正想要的候选完全不同",
            "传统输入法主要依赖静态词频和局部调频，缺少对上下文的理解",
            "用户经常翻页、选词、回退，这会明显打断表达节奏",
        ],
        size=18,
    )
    card(slide, Inches(0.95), Inches(4.15), Inches(3.6), Inches(1.5), "考试语境", ["输入：fx", "真正想要：复习", "而不是：分析"], accent=RED, title_size=18, body_size=13)
    card(slide, Inches(4.85), Inches(4.15), Inches(3.6), Inches(1.5), "发布语境", ["输入：bg", "真正想要：变更日志", "而不是：报告"], accent=CYAN, title_size=18, body_size=13)
    card(slide, Inches(8.75), Inches(4.15), Inches(3.6), Inches(1.5), "主持语境", ["输入：zc", "真正想要：主持词", "而不是：支持"], accent=GREEN, title_size=18, body_size=13)
    footer(slide)

    # 3 demo with pinyin
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band(slide, "02 演示一", CYAN)
    title_block(slide, "有拼音输入时，系统会根据语境重排候选", "这一页直接展示当前实测结果")
    headers = ["场景", "上下文", "输入", "原始候选池（未重排）", "重排后顺序"]
    rows = [
        ["学习备考", "马上要考试了，我需要开始……", "fx", "分析 / 复习 / 发现 / 风险 / 分享 / 方向", "复习 / 分析 / 分享 / 发现 / 方向 / 风险"],
        ["出差出行", "下周要去杭州出差，先把酒店和……", "jp", "精品 / 奖品 / 窘迫 / 键盘 / 绝配 / 机票", "机票 / 奖品 / 绝配 / 精品 / 键盘 / 窘迫"],
        ["技术开发", "今晚继续优化输入法 DLL 的……", "cp", "产品 / 测评 / 脆皮 / 藏品 / 参评 / 重排延迟", "重排延迟 / 产品 / 参评 / 藏品 / 测评 / 脆皮"],
    ]
    add_table_slide_rows(slide, Inches(0.72), Inches(1.9), Inches(11.9), Inches(3.95), headers, rows)
    callout(slide, Inches(0.9), Inches(6.05), Inches(11.55), Inches(0.55), "重点不是候选更多，而是在同样候选池里，系统能把更符合当前语境的结果提前。", fill=WHITE, line=LINE, color=SLATE, size=15)
    footer(slide)

    # 4 demo no input
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band(slide, "03 演示二", GREEN)
    title_block(slide, "用户停下来时，输入法还能继续帮助表达", "无输入预测不是替用户写全文，而是在停顿时降低继续表达的成本")
    card(slide, Inches(0.95), Inches(1.95), Inches(3.7), Inches(3.7), "技术语境", ["示例上下文：今晚继续优化输入法 DLL 的……", "展示系统补出的短结果", "强调它更像“继续表达”而不是“继续打字”"], accent=ORANGE, title_size=20, body_size=14)
    card(slide, Inches(4.82), Inches(1.95), Inches(3.7), Inches(3.7), "办公语境", ["示例上下文：先把发布说明和……", "展示系统补出的短结果", "说明输入法开始进入写作辅助"], accent=CYAN, title_size=20, body_size=14)
    card(slide, Inches(8.69), Inches(1.95), Inches(3.7), Inches(3.7), "学习语境", ["示例上下文：我打算先做几套……", "展示系统补出的短结果", "说明它能帮助用户继续组织内容"], accent=GREEN, title_size=20, body_size=14)
    footer(slide)

    # 5 architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band(slide, "04 系统方案", RED)
    title_block(slide, "系统方案：生成、重排、预测三层协作")
    arch_box(slide, Inches(0.8), Inches(2.1), Inches(3.1), Inches(1.8), "候选生成层", "Rime / 万象\n负责拼音输入、词库、基础候选生成", NAVY)
    arrow(slide, Inches(4.0), Inches(2.75))
    arch_box(slide, Inches(4.55), Inches(2.1), Inches(3.1), Inches(1.8), "候选重排层", "Alpha\n结合最近上下文，对候选池做实时语义重排", ORANGE)
    arrow(slide, Inches(7.75), Inches(2.75))
    arch_box(slide, Inches(8.3), Inches(2.1), Inches(4.1), Inches(1.8), "表达预测层", "LLM\n在无拼音场景下补出短词、短语或短句", CYAN, text_color=NAVY)
    callout(slide, Inches(1.0), Inches(4.65), Inches(11.2), Inches(0.75), "这不是把 AI 硬接到输入法后面，而是形成了职责清晰、可以持续扩展的结构化方案。", fill=RGBColor(239, 246, 255), line=CYAN, color=NAVY, size=16)
    bullets(
        slide,
        Inches(1.1),
        Inches(5.7),
        Inches(11.0),
        Inches(0.8),
        ["用户每一次上屏行为，又会成为后续排序和预测的反馈信号。"],
        size=15,
        color=SLATE,
    )
    footer(slide)

    # 6 innovation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band(slide, "05 核心创新", ORANGE)
    title_block(slide, "项目真正的创新，不在界面，而在交互范式")
    card(slide, Inches(0.9), Inches(1.95), Inches(3.8), Inches(1.65), "创新点 1", ["从“拼音匹配”升级到“语义重排”"], accent=ORANGE, title_size=18, body_size=14)
    card(slide, Inches(4.78), Inches(1.95), Inches(3.8), Inches(1.65), "创新点 2", ["输入法被拆成“生成层 + 排序层 + 预测层”"], accent=CYAN, title_size=18, body_size=14)
    card(slide, Inches(8.66), Inches(1.95), Inches(3.8), Inches(1.65), "创新点 3", ["用户上屏行为不再只是结果，而是未来排序的反馈信号"], accent=GREEN, title_size=18, body_size=14)
    card(slide, Inches(1.8), Inches(4.0), Inches(4.25), Inches(1.7), "创新点 4", ["强调本地实时智能，而不是纯云端依赖"], accent=RED, title_size=18, body_size=14)
    card(slide, Inches(7.0), Inches(4.0), Inches(4.25), Inches(1.7), "创新点 5", ["把输入法从“字符工具”推进到“表达协作系统”"], accent=NAVY, title_size=18, body_size=14)
    footer(slide)

    # 7 value
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band(slide, "06 价值与结尾", NAVY)
    title_block(slide, "为什么这个项目值得继续做")
    card(slide, Inches(0.95), Inches(1.95), Inches(3.65), Inches(2.0), "学习场景", ["真题", "复习", "知识点整理"], accent=ORANGE, title_size=20, body_size=15)
    card(slide, Inches(4.84), Inches(1.95), Inches(3.65), Inches(2.0), "办公场景", ["变更日志", "主持词", "发布说明"], accent=CYAN, title_size=20, body_size=15)
    card(slide, Inches(8.73), Inches(1.95), Inches(3.65), Inches(2.0), "技术场景", ["重排延迟", "缓存优化", "调试表达"], accent=GREEN, title_size=20, body_size=15)
    callout(slide, Inches(0.95), Inches(4.45), Inches(11.45), Inches(1.0), "Wisdom-Weasel 的目标不是单纯让用户“打字更快”，而是让用户“更顺畅、更低成本地完成表达”。", fill=RGBColor(255, 247, 237), line=AMBER, color=NAVY, size=18)
    callout(slide, Inches(0.95), Inches(5.8), Inches(11.45), Inches(0.72), "未来的输入法，不应该只是把拼音还原成汉字，而应该帮助人更顺畅地完成表达。", fill=WHITE, line=LINE, color=SLATE, size=16)
    footer(slide, "谢谢评委老师，欢迎提问。")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))


if __name__ == "__main__":
    build()
